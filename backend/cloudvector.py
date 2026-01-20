import os
import tempfile
import tiktoken
from pypdf import PdfReader
# from docx import Document
from pptx import Presentation
from supabase import create_client
from langchain_openai import OpenAIEmbeddings
from openai import OpenAI
from supabase.client import Client
from dotenv import load_dotenv

load_dotenv()

url: str = os.getenv("SUPABASE_URL")
key: str = os.getenv("SUPABASE_SERVICE_KEY")
supabase: Client = create_client(url, key)
 
client = OpenAI()
 
tokenizer = tiktoken.get_encoding("cl100k_base")
 
def download_from_supabase(file_path: str) -> str:
    """
    Downloads a file from Supabase Storage and returns a local file path.
    """
    bucket = "Documents"
    file_name = os.path.basename(file_path)
 
    tmp_dir = tempfile.gettempdir()
    local_path = os.path.join(tmp_dir, file_name)
 
    res = supabase.storage.from_(bucket).download(file_path)
 
    with open(local_path, "wb") as f:
        f.write(res)
 
    return local_path
 
def extract_text(file_path: str) -> str:

    local_path = download_from_supabase(file_path)

    if local_path.endswith(".pdf"):
        reader = PdfReader(local_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
 
    # if local_path.endswith(".docx"):
    #     doc = Document(local_path)
    #     return "\n".join(p.text for p in doc.paragraphs)
 
    if local_path.endswith(".pptx"):
        prs = Presentation(local_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
 
    raise ValueError("Unsupported file type")
 
def chunk_text(text, chunk_size=300, overlap=70):
    tokens = tokenizer.encode(text)
    chunks = []
 
    i = 0
    while i < len(tokens):
        chunk = tokens[i:i + chunk_size]
        chunks.append(tokenizer.decode(chunk))
        i += chunk_size - overlap
 
    return chunks
 
def embed_text(text: str) -> list[float]:
    response = client.embeddings.create(
        model="text-embedding-3-small",  # 1536 dims
        input=text
    )
    return response.data[0].embedding
 
def index_document(document_id: str, local_path: str):
    text = extract_text(local_path)
    chunks = chunk_text(text)
 
    # 1. Delete old chunks
    supabase.table("document_chunks") \
        .delete() \
        .eq("document_id", document_id) \
        .execute()
 
    # 2. Insert fresh chunks
    rows = []
    for i, chunk in enumerate(chunks):
        rows.append({
            "document_id": document_id,
            "chunk_index": i,
            "content": chunk,
            "embedding": embed_text(chunk),
        })
 
    r = supabase.table("document_chunks").insert(rows).execute()
 
def retrieve_chunks(query_embedding, limit=3):
    response = supabase.rpc(
        "match_embeddings",
        {
            "query_embedding": query_embedding,
            "match_count": limit
        }
    ).execute()
 
    return response.data
 
def format_citations(chunks):
   
    sources = []
 
    for c in chunks:
        page = c["page"]
        citation = {"chunk_id": c["chunk_id"], "text": c["content"], "source": c["file_name"], "page": f"{page}"}
        sources.append(citation)
 
    return sources
 
def chat(question: str):
    query_embedding = embed_text(question)
    chunks = retrieve_chunks(query_embedding)
 
    sources = format_citations(chunks)
 
    answer = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "system",
                "content": (
                    "Answer ONLY using the provided context. "
                )
            },
            {
                "role": "user",
                "content": f"Context:\n{sources}\n\nQuestion:\n{question}"
            }
        ]
    ).choices[0].message.content
 
    message = {"answer": answer, "question": question, "sources": sources }
 
    return message
 
def embed_all():
    files = (
        supabase.table("documents")
        .select("*")
        .execute()
    ).data
 
    for file in files:
        
        file_path = file["file_path"]
        file_id = file["id"]
        print(file_id, file_path)
 
        index_document(file_id, file_path)
 
# def test():
#     file_path = (
#         supabase.table("documents")
#         .select("file_path")
#         .eq("id", "2c6bf0de-f1de-49da-b4a0-e44caae7c344")
#         .execute()
#     ).data
#     print(file_path)
#     text = extract_text(file_path)
#     print(text)
#     index_document("2c6bf0de-f1de-49da-b4a0-e44caae7c344", file_path)

# print(chat("IBM Turbonomics")["answer"])
