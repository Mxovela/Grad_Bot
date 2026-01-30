import os
import tempfile
import tiktoken
from datetime import datetime
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openai import OpenAI
from dotenv import load_dotenv
from supabase_client import supabase
 
load_dotenv()
 
# Helps size text chunks based on token count
tokenizer = tiktoken.get_encoding("o200k_base")
 
# Embedding parameters
CHUNCK_SIZE = 180
CHUNCK_OVERLAP = 40

# Quering parameters
NO_PREVIOUS_QS = 3
LIMIT_HISTORY = 2*NO_PREVIOUS_QS # Must be even

# Sources parameters
RETRIEVAL_LIMIT = 10
QUERY_LIMIT = 5
CONTEXT_TOKEN_BUDGET = 2500
 
# Set up OpenAI client and parameters
client = OpenAI()

TEMP = 0.2
SIMILARITY_THRESHOLD = 0.1
PROMPT = (
        "You are a helpful assistant.\n"
        "Answer the user's question using ONLY the provided document context.\n"
        "You may expand on the answer using the provided context.\n"
        "If the answer is not present in the context, say you do not know.\n"
        "Strict rules:\n"
        "- Output valid Markdown only\n"
        "- Do NOT wrap the entire response in code blocks\n"
        "- Do NOT include explanations, notes, or metadata\n"
        "- Do NOT include markdown headers (#, ##, ###)\n"
        "Formatting rules:\n"
        "- Use **bold** for section titles and important terms\n"
        "- Use \"-\" for unordered lists\n"
        "- Use short paragraphs (1–2 lines max)\n"
        "- Leave a blank line between paragraphs and lists\n"
        "- Use inline code (`like this`) only for technical terms\n"
        "- Avoid nested lists\n"
        "Content rules:\n"
        "- Do NOT change meaning or add new information\n"
        "- Do NOT remove relevant information\n"
        "- Reformat only for readability in a chat bubble\n"
    )
 
TEST = True
 
## Train context
#/ Locally store and process the file to extract text
def download_from_supabase(file_path: str) -> list[dict]:
    """
    Downloads a file from Supabase Storage and returns a local file path.
    """
    bucket = "Documents"
    file_name = os.path.basename(file_path)
 
    tmp_dir = tempfile.gettempdir()
    local_path = os.path.join(tmp_dir, file_name)
 
    res = supabase.storage.from_(bucket).download(file_path)
 
    with open(local_path, "wb") as f:
        f.write(res)
 
    return local_path
 
def extract_pages(file_path: str) -> str:
 
    local_path = download_from_supabase(file_path)
 
    if local_path.endswith(".pdf"):
        reader = PdfReader(local_path)
        pages = []

        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append({
                    "page": i + 1,
                    "text": text
                })

        return pages
 
    if local_path.endswith(".docx"):
        doc = Document(local_path)
        return [{
            "page": 1,
            "text": "\n".join(p.text for p in doc.paragraphs)
        }]

    if local_path.endswith(".pptx"):
        prs = Presentation(local_path)
        pages = []

        for i, slide in enumerate(prs.slides):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)

            pages.append({
                "page": i + 1,
                "text": "\n".join(text)
            })

        return pages
 
    raise ValueError("Unsupported file type")
 
#/ Embed and upload: Split, embed and store text into chunks based on token count
def embed_text(text: str) -> list[float]:
    response = client.embeddings.create(
        model="text-embedding-3-small",  # 1536 dims
        input=text
    )
    return response.data[0].embedding
 
def chunk_pages(pages, chunk_size=CHUNCK_SIZE, overlap=CHUNCK_OVERLAP):
    chunks = []

    for page in pages:

        tokens = tokenizer.encode(page["text"])
        if len(tokens) < 10:
            continue

        i = 0
        while i < len(tokens):
            chunk_tokens = tokens[i:i + chunk_size]
            chunks.append({
                "text": tokenizer.decode(chunk_tokens),
                "page": page["page"]
            })
            i += chunk_size - overlap

    return chunks
 
def index_document(document_id: str, file_path: str):

    pages = extract_pages(file_path)
    chunks = chunk_pages(pages)

    supabase.table("document_chunks") \
        .delete() \
        .eq("document_id", document_id) \
        .execute()

    rows = []
    for i, chunk in enumerate(chunks):
        rows.append({
            "document_id": document_id,
            "chunk_index": i,
            "content": chunk["text"],
            "page": chunk["page"],
            "embedding": embed_text(chunk["text"]),
        })

    supabase.table("document_chunks").insert(rows).execute()
 
#/ Batch embedding
def embed_all():
    start_time = datetime.now().strftime('%I:%M:%S %p')
    print("Starting to embed all documents: ",start_time)
    files = (
        supabase.table("documents")
        .select("*")
        .execute()
    ).data
 
    for file in files:
       
        file_path = file["file_path"]
        file_id = file["id"]
 
        index_document(file_id, file_path)

    end_time = datetime.now().strftime('%I:%M:%S %p')
    print("Finished embedding all documents: ",end_time)
   
    return True
 
## Prompt
#/ Question
def question_history(chat_id):
 
    user_history_resp = (
        supabase.table("chat_messages")
        .select("content")
        .eq("chat_id", chat_id)
        .eq("role", "user")
        .order("id", desc=True)
        .limit(NO_PREVIOUS_QS)
        .execute()
    )
 
    previous_questions = [ r["content"] for r in user_history_resp.data]
 
    return previous_questions

def detect_query_intent(question: str) -> str:
    q = question.lower().strip()

    if q.startswith(("what is", "what’s", "what are", "define", "explain")):
        return "definition"

    if q.startswith(("how","how does", "how do", "how is", "how can")):
        return "mechanism"

    if q.startswith(("why",)):
        return "reasoning"

    if any(opt in q for opt in ["compare", "difference", "vs"]):
        return "comparison"

    return "general"

def rewrite_query(previous_questions, current_question) -> str:
 
    intent = detect_query_intent(current_question)

    if not previous_questions:
        return current_question

    history = "\n".join(previous_questions[::-1])

    if intent == "definition":
        instruction = (
            "Rewrite the user's question as a concise technical definition query.\n"
            "Include the product or concept category and its primary purpose.\n"
            "Do NOT include conversational phrasing.\n"
        )

    elif intent == "mechanism":
        instruction = (
            "Rewrite the user's question as a technical explanation query.\n"
            "Focus on how the system works or operates.\n"
        )

    elif intent == "comparison":
        instruction = (
            "Rewrite the user's question as a comparison query.\n"
            "Explicitly name the entities being compared.\n"
        )

    else:
        instruction = (
            "Rewrite the user's latest question into a standalone, explicit search query\n"
            "that can be used to retrieve relevant document passages.\n"
        )

    prompt = f"""
        Given the conversation history and the latest user question, {intent} intent,
        {instruction}

        Conversation (user questions only):
        {history}

        Latest question:
        {current_question}

        Standalone search query:
        """.strip()

    response = client.responses.create(
        model="gpt-4o-mini",
        input=prompt
    )

    return response.output_text.strip()
 
#/ Sources
def retrieve_chunks(query_embedding, limit=RETRIEVAL_LIMIT):
    response = supabase.rpc(
        "match_embeddings",
        {
            "query_embedding": query_embedding,
            "match_count": limit
        }
    ).execute()

    if not response.data:
        return []

    # Filter by similarity threshold
    filtered = []
    for r in response.data:

        print("file: ",r["file_name"]," similarity: ",r["similarity"])

        if r["similarity"] >= SIMILARITY_THRESHOLD:
            filtered.append(r)

    if not filtered: filtered = [r for r in response.data if r["similarity"] >= 0]  # Ensure at least one chunk is returned
    if not filtered: filtered = [response.data[0]]  # Fallback to at least one chunk if none meet criteria

    if len(filtered) > QUERY_LIMIT:
        filtered = filtered[:QUERY_LIMIT]

    return filtered
 
def trim_chunks_by_tokens(chunks: list[dict]) -> list[dict]:
    total_tokens = 0
    sources = []
 
    for chunk in chunks:
 
        tokens = len(tokenizer.encode(chunk["content"]))
        if total_tokens + tokens > CONTEXT_TOKEN_BUDGET:
            break
 
        sources.append(format_citations(chunk))
 
        total_tokens += tokens
 
    return sources
 
def format_citations(chunk):
   
    page = chunk["page"]
    citation = {"chunk_id": chunk["chunk_id"], "text": chunk["content"], "source": chunk["file_name"], "page": f"{page}"}
 
    return citation

def chunks_by_id(chunks: list[str]) -> list[dict]:
    sources = []

    for chunk in chunks:
        response = supabase.rpc(
            "citations",
            {
                "p_chunk_id": chunk,
            }
        ).execute()

        for source in response.data:
            sources.append(format_citations(source))

    return sources


#/ Chat history
def ordered_history(user_id, limit=100):
 
    chat_id = get_chat_id(user_id)
   
    if not chat_id:
        chat_id = new_chat(user_id)
 
    chat_history_resp = (
        supabase.table("chat_messages")
        .select("*")
        .eq("chat_id", chat_id)
        .order("id", desc=True)
        .limit(limit)
        .execute()
    )
 
    chat_history = list(reversed(chat_history_resp.data))
 
    return chat_history
 
def context_history(chat_id, limit=LIMIT_HISTORY):
    chat_history_resp = (
        supabase.table("chat_messages")
        .select("role, content")
        .eq("chat_id", chat_id)
        .order("id", desc=True)
        .limit(limit)
        .execute()
    )
 
    chat_history = list(reversed(chat_history_resp.data))
 
    return chat_history
 
def messages(chat_history, context_text, question):
    messages = [{"role": "system", "content": PROMPT}]
 
    for msg in chat_history:
        messages.append({
            "role": msg["role"],
            "content": msg["content"]
        })
   
    messages.append({
        "role": "user",
        "content": f"""
        Context:
        {context_text}
 
        Question:
        {question}
        """.strip()
            })
   
    return messages

## Chat Handling
def get_chat_id(user_id):
    chat_id = (
        supabase.table("chats")
        .select("id")
        .eq("user_id", user_id)
        .execute()
        ).data
    
    if not chat_id:
        return None
   
    return chat_id[0]["id"]
 
def new_chat(user_id):
    chat_id = (
        supabase.table("chats")
        .insert({"user_id": user_id})
        .execute()
        ).data
    
    if not chat_id:
        return None
   
    return chat_id[0]["id"]
 
def delete_chat(chat_id):
    response = (
        supabase.table("chats")
        .delete()
        .eq("id", chat_id)
        .execute()
        )
 
def new_chat_on_token(user_id):
    old_chat = get_chat_id(user_id)
    if old_chat:
        delete_chat(old_chat)
 
    chat_id = new_chat(user_id)
    return chat_id
 
def get_user_message_count(user_id):
    chat_id = get_chat_id(user_id)
    if not chat_id:
        return 0
    
    count = (
        supabase.table("chat_messages")
        .select("*", count="exact")
        .eq("chat_id", chat_id)
        .eq("role", "user")
        .execute()
    ).count
    
    return count if count else 0

## Chat function to answer questions based on document context
def chat(user_id, question):
 
    local_time = datetime.now().strftime('%I:%M:%S %p')
 
    chat_id = get_chat_id(user_id)
 
    if not chat_id:
        chat_id = new_chat(user_id)

    if not chat_id:
        return {"error": "Could not create chat session."}
 
    ##QUESTIONS
    #/ improve question
    previous_questions = question_history(chat_id)

    if not previous_questions: previous_questions = [question]
 
    rewritten_question = rewrite_query(previous_questions, question)
    if TEST:
        print("previous_questions: ","\n".join(previous_questions))
        print("question:           ",question)
        print("rewritten_question: ",rewritten_question)
   
    #/ embed question
    query_embedding = embed_text(rewritten_question)
 
    ##SOURCES
    #/ get the best sources
    chunks = retrieve_chunks(query_embedding)
    sources = trim_chunks_by_tokens(chunks)
 
    #/ format context for the LLM
    context_text = ''
    for source in sources:
        file = source["source"]
        page = source["page"]
        content = source["text"]
        context_text = context_text + f"Source: {file}, Page:{page} \n {content} \n \n"
 
    #/ get chat history
    chat_history = context_history(chat_id)
    
    #/ build messages
    input = messages(chat_history, context_text, question)
 
    # if TEST: 
    #     print("input: ", (input))
 
    ##Answer
    #/ get response
    completion = client.responses.create(
        model="gpt-4o-mini",
        temperature=TEMP,
        input=input
    )

    answer = completion.output_text.strip()
 
    if TEST: print("\n answer: ",answer)

    ## Update DB

    #/ store question
    question_tokens = len(tokenizer.encode(question))

    q = supabase.table("chat_messages").insert({
        "chat_id": chat_id,
        "role": "user",
        "content": question,
        "token_count": question_tokens,
        "time_stamp": local_time
    }).execute()
 
    #/ store answer
    answer_time = datetime.now().strftime('%I:%M:%S %p')
    answer_tokens = len(tokenizer.encode(answer))
    sources_id = [ s["chunk_id"] for s in sources]
 
    supabase.table("chat_messages").insert({
        "chat_id": chat_id,
        "role": "assistant",
        "content": answer,
        "time_stamp": answer_time,
        "token_count": answer_tokens,
        "sources": sources_id
        }).execute()
 
    message = {"answer": answer, "question": question, "sources": sources, "chat_id": chat_id}
 
    return message

def test():
    #embed_all()
    # user_id = "73205ea6-2afe-4407-a35e-6ea6f7260333"
    question = "What is instana?"
    convo  = chat(user_id, question)
    q = convo["question"]
    a = convo["answer"]
    print("\n number of sources:", len(convo["sources"]))

    prompt = f"""
        A bot is used to answer trainees' questions based on provided document context.
        Rate the question and answer pair out of 100. give a short explanation. Briefly suggest a better question if necessary and ways to improve the bot.
        structure it as:
          rating: <score out of 100>
          explanation: <short explanation>
          suggested_bot_improvements: <suggestions>
          suggested_question_improvements: <suggestions>
 
        Question: {q}
        Answer: {a}
        """.strip()
 
    response = client.responses.create(
        model="gpt-4o-mini",
        input=prompt
    )

    print("\n",response.output_text.strip())

