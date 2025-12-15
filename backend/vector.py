from dotenv import load_dotenv
from langchain_openai import OpenAIEmbeddings
from langchain_chroma import Chroma
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
import os
import hashlib

load_dotenv()

# ==============================
# DOCUMENT LOADING
# ==============================

def load_documents_from_folder(folder_path: str = "documents") -> list[Document]:
    documents: list[Document] = []

    # Optional imports
    try:
        from pypdf import PdfReader
    except ImportError:
        PdfReader = None

    try:
        from docx import Document as DocxDocument
    except ImportError:
        DocxDocument = None

    try:
        from pptx import Presentation
    except ImportError:
        Presentation = None

    supported_extensions = {".pdf", ".docx", ".pptx", ".txt", ".md"}

    for root, _, files in os.walk(folder_path):
        for file in files:
            file_ext = os.path.splitext(file)[1].lower()
            if file_ext not in supported_extensions:
                continue

            file_path = os.path.join(root, file)
            relative_path = os.path.relpath(file_path, folder_path)
            file_name = os.path.basename(file_path)

            try:
                # ==============================
                # PDF (PAGE-AWARE)
                # ==============================
                if file_ext == ".pdf":
                    if PdfReader is None:
                        continue

                    reader = PdfReader(file_path)

                    for page_number, page in enumerate(reader.pages, start=1):
                        page_text = page.extract_text()
                        if not page_text or not page_text.strip():
                            continue

                        documents.append(
                            Document(
                                page_content=page_text,
                                metadata={
                                    "source": file_name,
                                    "page": page_number,
                                    "file_path": relative_path,
                                    "file_type": "pdf"
                                }
                            )
                        )

                # ==============================
                # DOCX
                # ==============================
                elif file_ext == ".docx":
                    if DocxDocument is None:
                        continue

                    doc = DocxDocument(file_path)
                    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

                    if text:
                        documents.append(
                            Document(
                                page_content=text,
                                metadata={
                                    "source": file_name,
                                    "page": None,
                                    "file_path": relative_path,
                                    "file_type": "docx"
                                }
                            )
                        )

                # ==============================
                # PPTX
                # ==============================
                elif file_ext == ".pptx":
                    if Presentation is None:
                        continue

                    prs = Presentation(file_path)
                    slide_texts = []

                    for slide_idx, slide in enumerate(prs.slides, start=1):
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_texts.append(f"Slide {slide_idx}: {shape.text}")

                    if slide_texts:
                        documents.append(
                            Document(
                                page_content="\n".join(slide_texts),
                                metadata={
                                    "source": file_name,
                                    "page": None,
                                    "file_path": relative_path,
                                    "file_type": "pptx"
                                }
                            )
                        )

                # ==============================
                # TXT / MD
                # ==============================
                elif file_ext in {".txt", ".md"}:
                    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                        text = f.read()

                    if text.strip():
                        documents.append(
                            Document(
                                page_content=text,
                                metadata={
                                    "source": file_name,
                                    "page": None,
                                    "file_path": relative_path,
                                    "file_type": file_ext[1:]
                                }
                            )
                        )

            except Exception as e:
                print(f"Error processing {file_path}: {e}")

    return documents


# ==============================
# CHUNKING
# ==============================

def chunk_documents(documents: list[Document]) -> list[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200
    )

    chunked_docs = splitter.split_documents(documents)

    for idx, doc in enumerate(chunked_docs):
        doc.metadata["chunk_id"] = f"chunk_{idx}"

        # Stable unique ID per chunk
        doc.id = hashlib.md5(
            f"{doc.metadata.get('source')}-{doc.metadata.get('page')}-{idx}".encode()
        ).hexdigest()

    return chunked_docs


# ==============================
# VECTOR STORE SETUP
# ==============================

embeddings = OpenAIEmbeddings(model="text-embedding-3-small")

DB_PATH = "./chrome_langchain_db"
add_documents = not os.path.exists(DB_PATH)

vector_store = Chroma(
    collection_name="grad_knowledge",
    persist_directory=DB_PATH,
    embedding_function=embeddings
)

if add_documents:
    raw_documents = load_documents_from_folder("documents")
    documents = chunk_documents(raw_documents)
    ids = [doc.id for doc in documents]

    vector_store.add_documents(documents=documents, ids=ids)
    vector_store.persist()

retriever = vector_store.as_retriever(
    search_kwargs={"k": 5}
)
